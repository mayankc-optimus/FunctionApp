import os
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from io import BytesIO
import azure.functions as func
import logging

app = func.FunctionApp()

def read_file(file_path):
    ext = file_path.lower().split('.')[-1]
    if not os.path.exists(file_path):
        return f"Error: The file at {file_path} does not exist."

    try:
        with open(file_path, 'rb') as file:
            file_content = file.read()

            # For Text Files (.txt)
            if ext == 'txt':
                return file_content.decode('utf-8')
            
            # For PDF Files (.pdf)
            elif ext == 'pdf':
                reader = PyPDF2.PdfReader(BytesIO(file_content))
                text = ''
                for page in reader.pages:
                    text += page.extract_text()
                return text
            
            # For Excel Files (.xls, .xlsx)
            elif ext in ['xlsx', 'xls']:
                df = pd.read_excel(BytesIO(file_content), sheet_name=0)
                return df.to_string(index=False)

            # For PowerPoint Files (.pptx)
            elif ext == 'pptx':
                prs = Presentation(BytesIO(file_content))
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + '\n'
                return text

            # For Word Documents (.docx)
            elif ext == 'docx':
                doc = Document(BytesIO(file_content))
                text = '\n'.join([para.text for para in doc.paragraphs])
                return text

            else:
                return "Unsupported file type."

    except Exception as e:
        return f"Error reading file: {str(e)}"

@app.route(route="ReadFilesFunction", auth_level=func.AuthLevel.Anonymous)
def ReadFilesFunction(req: func.HttpRequest) -> func.HttpResponse:
    logging.info('Python HTTP trigger function processed a request.')

    try:
        # Get the file path from the query parameters
        file_path = req.params.get('file_path')
        if not file_path:
            return func.HttpResponse("Please pass the 'file_path' in the query string.", status_code=400)

        # Process the file
        file_content = read_file(file_path)

        # Return the result
        return func.HttpResponse(file_content, mimetype="text/plain")
    
    except Exception as e:
        return func.HttpResponse(f"Error: {str(e)}", status_code=500)
